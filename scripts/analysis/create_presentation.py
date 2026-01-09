"""
生成研究进展汇报PPT - 中文版
基于研究报告创建专业演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

def create_presentation():
    """创建研究进展汇报PPT"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 幻灯片1 - 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 添加深蓝色背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 35, 126)  # 深蓝色

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "脑遮挡检测研究进展汇报"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "人类与AI视觉识别策略对比研究"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(255, 193, 7)  # 金色
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 日期
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "2026年1月9日"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = RGBColor(200, 200, 200)
    date_para.alignment = PP_ALIGN.CENTER

    # 幻灯片2 - 研究概述
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "研究背景与目标"
    title.text_frame.paragraphs[0].font.size = Pt(40)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "数据集：OIID遮挡飞机图像数据集"
    tf.paragraphs[0].font.size = Pt(24)

    for text in [
        "任务：二分类（Aircraft1 vs Aircraft2）",
        "遮挡级别：10%, 70%, 90%",
        "对比模型：ViT-B/16 vs ResNet-50",
        "核心问题：AI能否像人类一样处理遮挡？"
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(24)

    # 幻灯片3 - 核心发现
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "五大核心发现"
    title.text_frame.paragraphs[0].font.size = Pt(40)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. 巨大的人机差距：AI比人类差5-52%"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True

    for i, text in enumerate([
        "2. 架构很重要：ViT比ResNet好6%",
        "3. 反直觉现象：低遮挡时AI表现更差",
        "4. 小样本挑战：仅300张训练图像",
        "5. 数据增强尝试：扩展到9,900张图像"
    ], start=2):
        p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(24)
        p.font.bold = True

    # 幻灯片4 - 性能对比表格
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "人类 vs AI 性能对比"
    title.text_frame.paragraphs[0].font.size = Pt(40)

    # 添加表格
    rows, cols = 4, 6
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(3.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 表头
    headers = ["遮挡级别", "人类", "ViT", "ResNet", "ViT差距", "ResNet差距"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(18)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(26, 35, 126)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 数据
    data = [
        ["10%", "95.62%", "50.00%", "43.75%", "45.62%", "51.87%"],
        ["70%", "79.28%", "50.00%", "43.75%", "29.28%", "35.53%"],
        ["90%", "61.88%", "56.25%", "50.00%", "5.63%", "11.88%"]
    ]

    for i, row_data in enumerate(data, start=1):
        for j, value in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            # 差距列用红色标注
            if j >= 4:
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                cell.text_frame.paragraphs[0].font.bold = True

    # 幻灯片5 - 关键洞察
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "性能分析关键洞察"
    title.text_frame.paragraphs[0].font.size = Pt(40)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "低遮挡（10%）差距最大：45-52%"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "→ 假设：小样本过拟合"
    p.level = 1
    p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "高遮挡（90%）差距缩小：5-12%"
    p.level = 0
    p.font.size = Pt(24)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "→ 任务本身变得困难"
    p.level = 1
    p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "ViT全面优于ResNet"
    p.level = 0
    p.font.size = Pt(24)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "→ 全局注意力 > 局部卷积"
    p.level = 1
    p.font.size = Pt(22)

    # 幻灯片6 - 训练动态
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "训练动态分析"
    title.text_frame.paragraphs[0].font.size = Pt(40)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "ViT-B/16："
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    for text in [
        "最佳验证准确率：64.29%（第5轮）",
        "训练准确率：54.29%",
        "结论：欠拟合（训练<验证）"
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 1
        p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "ResNet-50："
    p.level = 0
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(231, 76, 60)

    for text in [
        "最佳验证准确率：54.76%（第4轮）",
        "训练准确率：47.62%",
        "结论：欠拟合，性能更差"
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 1
        p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "💡 启示：冻结主干网络策略过于保守"
    p.level = 0
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 193, 7)

    # 幻灯片7 - 当前挑战
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "面临的主要挑战"
    title.text_frame.paragraphs[0].font.size = Pt(40)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. 小样本问题：300张图像不足"
    tf.paragraphs[0].font.size = Pt(26)

    for i, text in enumerate([
        "2. 冻结主干限制：模型容量未充分利用",
        "3. 缺乏可解释性：不知道模型关注什么",
        "4. 架构探索不完整：仅测试2种模型"
    ], start=2):
        p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(26)

    # 幻灯片8 - 研究方向1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "方向1：遮挡感知注意力机制 ⭐⭐⭐"
    title.text_frame.paragraphs[0].font.size = Pt(36)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "动机：人类主动忽略遮挡区域"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "方法："
    p.level = 0
    p.font.size = Pt(24)
    p.font.bold = True

    for text in [
        "设计遮挡检测模块",
        "修改注意力机制降低遮挡区域权重",
        "增强可见区域特征提取"
    ]:
        p = tf.add_paragraph()
        p.text = "• " + text
        p.level = 1
        p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "预期效果：准确率提升10-20%"
    p.level = 0
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    # 幻灯片9 - 研究方向2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "方向2：基于部件的识别系统 ⭐⭐⭐"
    title.text_frame.paragraphs[0].font.size = Pt(36)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "动机：人类通过部件识别物体"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "（机翼、机身、尾翼）"
    p.level = 1
    p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text = "方法："
    p.level = 0
    p.font.size = Pt(24)
    p.font.bold = True

    for text in [
        "预训练部件检测器",
        "构建部件关系图",
        "使用图神经网络推理"
    ]:
        p = tf.add_paragraph()
        p.text = "• " + text
        p.level = 1
        p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "优势：鲁棒、可解释、符合人类认知"
    p.level = 0
    p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "预期效果：准确率提升15-25%"
    p.level = 0
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    # 幻灯片10 - 研究方向3
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "方向3：fMRI引导的模型设计 ⭐⭐⭐⭐"
    title.text_frame.paragraphs[0].font.size = Pt(36)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "动机：用人脑数据指导AI架构"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = "方法："
    p.level = 0
    p.font.size = Pt(24)
    p.font.bold = True

    for text in [
        "表征相似性分析（RSA）",
        "编码模型：AI特征预测fMRI信号",
        "脑启发架构设计"
    ]:
        p = tf.add_paragraph()
        p.text = "• " + text
        p.level = 1
        p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "优势："
    p.level = 0
    p.font.size = Pt(24)
    p.font.bold = True

    for text in [
        "理论基础强",
        "发表潜力高（Nature级别）",
        "连接AI与神经科学"
    ]:
        p = tf.add_paragraph()
        p.text = "✓ " + text
        p.level = 1
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0, 128, 0)

    p = tf.add_paragraph()
    p.text = "⚠ 挑战：需要fMRI数据"
    p.level = 0
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(255, 0, 0)

    # 幻灯片11 - 下一步计划
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "下一步行动计划"
    title.text_frame.paragraphs[0].font.size = Pt(40)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "近期（1-2周）："
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)

    for text in [
        "完成数据增强实验（9,900张）",
        "全模型微调（解冻主干）",
        "生成注意力可视化"
    ]:
        p = tf.add_paragraph()
        p.text = "✓ " + text
        p.level = 1
        p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "短期（1-2月）："
    p.level = 0
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 204, 113)

    for text in [
        "选项A：实现遮挡感知注意力（工程导向）",
        "选项B：fMRI验证分析（理论导向）"
    ]:
        p = tf.add_paragraph()
        p.text = "• " + text
        p.level = 1
        p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "中期（3-6月）："
    p.level = 0
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(155, 89, 182)

    for text in [
        "实现基于部件的识别",
        "扩展架构对比",
        "准备顶会论文"
    ]:
        p = tf.add_paragraph()
        p.text = "• " + text
        p.level = 1
        p.font.size = Pt(22)

    # 幻灯片12 - 关键决策
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "需要决策的关键问题"
    title.text_frame.paragraphs[0].font.size = Pt(40)

    content = slide.placeholders[1]
    tf = content.text_frame

    questions = [
        "1. 研究方向：工程（快速结果）vs 理论（高影响）？",
        "2. fMRI数据：OIID数据集中是否可用？",
        "3. 发表目标：会议（CVPR/NeurIPS）vs 期刊（Nature）？",
        "4. 时间线：第一篇论文的截止日期？",
        "5. 计算资源：是否有GPU集群支持？"
    ]

    tf.text = questions[0]
    tf.paragraphs[0].font.size = Pt(24)

    for q in questions[1:]:
        p = tf.add_paragraph()
        p.text = q
        p.level = 0
        p.font.size = Pt(24)

    # 幻灯片13 - 总结
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "研究总结"
    title.text_frame.paragraphs[0].font.size = Pt(40)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "已完成："
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)

    for text in [
        "Phase 1基线实验（ViT vs ResNet）",
        "量化人机性能差距（最大51.87%）",
        "发现架构差异（ViT优于ResNet）",
        "建立完整分析流程"
    ]:
        p = tf.add_paragraph()
        p.text = "✓ " + text
        p.level = 1
        p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "核心贡献："
    p.level = 0
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 193, 7)

    for text in [
        "揭示AI在遮挡场景下的巨大差距",
        "证明全局注意力优于局部卷积",
        "提出5个可行的研究方向"
    ]:
        p = tf.add_paragraph()
        p.text = "★ " + text
        p.level = 1
        p.font.size = Pt(22)

    # 幻灯片14 - 致谢
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加深蓝色背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 35, 126)

    # 谢谢标题
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "谢谢！"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(54)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = RGBColor(255, 255, 255)
    thanks_para.alignment = PP_ALIGN.CENTER

    # 欢迎提问
    question_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    question_frame = question_box.text_frame
    question_frame.text = "欢迎提问与建议"
    question_para = question_frame.paragraphs[0]
    question_para.font.size = Pt(28)
    question_para.font.color.rgb = RGBColor(255, 193, 7)
    question_para.alignment = PP_ALIGN.CENTER

    # 报告详情
    detail_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    detail_frame = detail_box.text_frame
    detail_frame.text = "报告详情：reports/research_analysis_20260109_223853.md\n可视化图表：reports/analysis_outputs/"
    detail_para = detail_frame.paragraphs[0]
    detail_para.font.size = Pt(16)
    detail_para.font.color.rgb = RGBColor(200, 200, 200)
    detail_para.alignment = PP_ALIGN.CENTER

    # 保存PPT
    output_path = Path("D:/University/Junior/1st/code/brain_occlusion_detection/reports/研究进展汇报_20260109.pptx")
    prs.save(str(output_path))

    print(f"[OK] PPT已生成：{output_path}")
    print(f"\nPPT包含：")
    print("  - 14张幻灯片")
    print("  - 封面（深蓝色主题）")
    print("  - 研究概述、核心发现")
    print("  - 性能对比表格（彩色标注）")
    print("  - 训练动态分析")
    print("  - 3个主要研究方向（详细说明）")
    print("  - 行动计划（近期/短期/中期）")
    print("  - 关键决策问题")
    print("  - 研究总结")
    print("  - 致谢页面")
    print("\n设计特点：")
    print("  - 专业学术风格")
    print("  - 深蓝色+金色配色")
    print("  - 大字体（24-40pt）")
    print("  - 清晰的层次结构")
    print("  - 适合15分钟汇报")

    return output_path

if __name__ == "__main__":
    create_presentation()
